"""
pdf2ppt_optimized workflow
~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
Optimized version of pdf2ppt_parallel:
1.  **Text Extraction**: Uses MinerU blocks for both layout and text content, relying on its layout analysis
    to ensure correct column separation.
2.  **Dynamic Slide Sizing**: Sets PPT slide dimensions to match input image resolution exactly.
3.  **Smart Text Styling**: 
    -   Samples text color from the original image.
    -   Auto-centers titles.
    -   Calculates optimal font size.
"""

from __future__ import annotations
import os
import asyncio
from pathlib import Path
from typing import List, Dict, Any, Optional
from collections import Counter
import copy
import time
import math

import cv2
import numpy as np
import fitz  # PyMuPDF
import yaml
from PIL import Image

from dataflow_agent.workflow.registry import register
from dataflow_agent.graphbuilder.graph_builder import GenericGraphBuilder
from dataflow_agent.logger import get_logger

from dataflow_agent.state import Paper2FigureState
from dataflow_agent.utils import get_project_root, pixels_to_inches, calculate_font_size

# Tools
from dataflow_agent.toolkits.multimodaltool.sam_tool import segment_layout_boxes, segment_layout_boxes_server, free_sam_model
from dataflow_agent.toolkits.multimodaltool.bg_tool import local_tool_for_bg_remove, free_bg_rm_model
from dataflow_agent.toolkits.multimodaltool.mineru_tool import recursive_mineru_layout
from dataflow_agent.toolkits.multimodaltool.req_img import generate_or_edit_and_save_image_async
from dataflow_agent.toolkits.multimodaltool import ppt_tool

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

from dataflow_agent.toolkits.multimodaltool.ppt_text_fit import DEFAULT_FITTER, TextFitStyle

log = get_logger(__name__)

# Load configuration from yaml
def load_server_config():
    root = get_project_root()
    config_path = root / "conf" / "model_servers.yaml"
    if not config_path.exists():
        log.warning(f"Config file not found at {config_path}, using defaults.")
        return {}
    try:
        with open(config_path, "r") as f:
            return yaml.safe_load(f) or {}
    except Exception as e:
        log.error(f"Failed to load config: {e}")
        return {}

SERVER_CONFIG = load_server_config()

def _resize_image_for_ppt(img_path: str, max_w: int = 1920, max_h: int = 1080) -> str:
    """
    Force resize image to fit within max_w x max_h (contain), keeping aspect ratio.
    Overwrites the original image or saves a new one.
    This prevents huge images from crashing pptx (max 56 inches) and stabilizes MinerU/SAM results.
    """
    if not os.path.exists(img_path):
        return img_path
    
    try:
        with Image.open(img_path) as img:
            w, h = img.size
            # If image is already smaller than target box, skip (or optionally upscale?)
            # Here we only downscale to avoid quality loss on small icons, 
            # unless it's way too small? For now, mainly prevent huge images.
            if w <= max_w and h <= max_h:
                return img_path
            
            # Calculate scale to contain
            scale = min(max_w / w, max_h / h)
            new_w = int(w * scale)
            new_h = int(h * scale)
            
            log.info(f"[pdf2ppt_opt] Resizing input image from {w}x{h} to {new_w}x{new_h} (fit {max_w}x{max_h})")
            
            resized = img.resize((new_w, new_h), Image.Resampling.LANCZOS)
            
            # Save to a new path to keep original safe (optional, but good practice)
            # or overwrite if we want consistency for downstream tools.
            # Let's overwrite/update filename to indicate resized.
            p = Path(img_path)
            new_path = p.parent / f"{p.stem}_resized{p.suffix}"
            resized.save(new_path)
            return str(new_path)
            
    except Exception as e:
        log.error(f"[pdf2ppt_opt] Resize failed: {e}")
        return img_path

# Helper to construct URLs
def get_sam_urls():
    if os.environ.get("SAM_SERVER_URLS"):
        return os.environ.get("SAM_SERVER_URLS").split(",")
    sam_cfg = SERVER_CONFIG.get("sam", {})
    instances = sam_cfg.get("instances", [])
    if instances:
        urls = []
        for inst in instances:
            for port in inst.get("ports", []):
                urls.append(f"http://127.0.0.1:{port}")
        if urls:
            return urls
    return ["http://localhost:8021", "http://localhost:8022","http://localhost:8023"]

SAM_SERVER_URLS = get_sam_urls()

def _ensure_result_path(state: Paper2FigureState) -> str:
    raw = getattr(state, "result_path", None)
    if raw:
        return raw
    root = get_project_root()
    ts = int(__import__("time").time())
    base_dir = (root / "outputs" / "pdf2ppt_optimized" / str(ts)).resolve()
    base_dir.mkdir(parents=True, exist_ok=True)
    state.result_path = str(base_dir)
    return state.result_path

def _run_sam_on_pages(image_paths: List[str], base_dir: str) -> List[Dict[str, Any]]:
    results: List[Dict[str, Any]] = []
    sam_ckpt = f"{get_project_root()}/sam_b.pt"

    for page_idx, img_path in enumerate(image_paths):
        img_path_obj = Path(img_path)
        if not img_path_obj.exists():
            log.warning(f"[pdf2ppt_opt] image not found for SAM: {img_path}")
            results.append({"page_idx": page_idx, "layout_items": []})
            continue

        out_dir = Path(base_dir) / "layout_items" / f"page_{page_idx+1:03d}"
        out_dir.mkdir(parents=True, exist_ok=True)

        try:
            layout_items = segment_layout_boxes_server(
                image_path=str(img_path_obj),
                output_dir=str(out_dir),
                server_urls=SAM_SERVER_URLS,
                checkpoint=sam_ckpt,
                min_area=200,
                min_score=0.0,
                iou_threshold=0.4,
                top_k=25,
                nms_by="mask",
            )
        except Exception as e:
            log.error(f"[pdf2ppt_opt] Remote SAM failed: {e}. Fallback to local.")
            layout_items = segment_layout_boxes(
                image_path=str(img_path_obj),
                output_dir=str(out_dir),
                checkpoint=sam_ckpt,
                min_area=200,
                min_score=0.0,
                iou_threshold=0.4,
                top_k=25,
                nms_by="mask",
            )
            
        try:
            pil_img = Image.open(str(img_path_obj))
            w, h = pil_img.size
        except Exception as e:
            log.error(f"[pdf2ppt_opt][page#{page_idx+1}] open image failed: {e}")
            w, h = 1024, 768

        for it in layout_items:
            bbox = it.get("bbox")
            if bbox and len(bbox) == 4:
                x1n, y1n, x2n, y2n = bbox
                x1 = int(round(x1n * w))
                y1 = int(round(y1n * h))
                x2 = int(round(x2n * w))
                y2 = int(round(y2n * h))
                if x2 > x1 and y2 > y1:
                    it["bbox_px"] = [x1, y1, x2, y2]

        results.append({"page_idx": page_idx, "layout_items": layout_items})

    try:
        free_sam_model(checkpoint=sam_ckpt)
    except Exception as e:
        log.error(f"[pdf2ppt_opt] free_sam_model failed: {e}")

    return results

# ==============================================================================
# Helper Functions for PPT Generation
# ==============================================================================

def _get_dominant_color(img_path: str, bbox: List[int]) -> RGBColor:
    """
    Smartly extract text color from the image region.
    Logic:
    1. Identify background color (most frequent).
    2. Find text color (frequent color with high contrast to background).
    3. Fallback to Black/White based on background brightness.
    """
    try:
        with Image.open(img_path) as img:
            x1, y1, x2, y2 = bbox
            # Clamp
            x1 = max(0, x1)
            y1 = max(0, y1)
            x2 = min(img.width, x2)
            y2 = min(img.height, y2)
            
            if x2 <= x1 or y2 <= y1:
                return RGBColor(0, 0, 0)
            
            crop = img.crop((x1, y1, x2, y2))
            
            # Reduce resolution to speed up
            crop.thumbnail((50, 50))
            crop = crop.convert("RGB")
            
            # Get colors (count, (r,g,b))
            colors = crop.getcolors(maxcolors=2500)
            if not colors:
                return RGBColor(0, 0, 0)
            
            # Sort by frequency (descending)
            # Assumption: The most frequent color is the background
            colors.sort(key=lambda x: x[0], reverse=True)
            
            bg_rgb = colors[0][1]
            
            def get_brightness(c):
                return 0.299 * c[0] + 0.587 * c[1] + 0.114 * c[2]
            
            def get_dist(c1, c2):
                return abs(c1[0]-c2[0]) + abs(c1[1]-c2[1]) + abs(c1[2]-c2[2])

            bg_brightness = get_brightness(bg_rgb)
            
            # Find the first color that has enough contrast with background
            # We skip the first one (background itself)
            best_rgb = None
            
            for i in range(1, len(colors)):
                count, rgb = colors[i]
                # Ignore very rare artifacts
                if count < colors[0][0] * 0.05: 
                    continue
                    
                # Check contrast
                # 1. Brightness diff > 50
                # 2. Or total Euclidean-ish dist > 100
                bri = get_brightness(rgb)
                dist = get_dist(rgb, bg_rgb)
                
                if abs(bri - bg_brightness) > 60 or dist > 150:
                    best_rgb = rgb
                    break
            
            # Fallback if no contrasting color found (e.g. solid block)
            if best_rgb is None:
                # If background is dark, use white; else black
                if bg_brightness < 128:
                    best_rgb = (255, 255, 255)
                else:
                    best_rgb = (0, 0, 0)
            
            return RGBColor(best_rgb[0], best_rgb[1], best_rgb[2])
            
    except Exception as e:
        log.warning(f"[pdf2ppt_opt] Color sampling failed: {e}")
        return RGBColor(0, 0, 0)

def _estimate_font_size_from_slide_geometry_pt(
    text: str,
    bbox_px: List[int],
    text_level: int,
    *,
    dpi: int = 96,
    line_spacing: float = 1.0,
) -> int:
    """
    Stable font-size estimation based on *PPT physical size*, not raw pixels.

    We map bbox height (inches) -> points and then shrink by estimated wrapping.

    This fixes the common failure modes:
    - FIGURE mode: image is scaled into a fixed 16:9 canvas, bbox_px becomes smaller,
      but the PPT physical size stays consistent. Using inches keeps font size stable.
    - Huge input images: avoids tiny fonts after downscaling.
    - Small images: avoids huge fonts that overflow.
    """
    text = (text or "").strip()
    if not text:
        return 12

    x1, y1, x2, y2 = bbox_px
    box_w_px = max(1, x2 - x1)
    box_h_px = max(1, y2 - y1)

    box_w_in = pixels_to_inches(box_w_px, dpi=dpi)
    box_h_in = pixels_to_inches(box_h_px, dpi=dpi)
    box_w_pt = box_w_in * 72.0
    box_h_pt = box_h_in * 72.0

    # Base ratio by level (tunable)
    if text_level == 1:  # title/header
        ratio = 0.75
        min_pt, max_pt = 14, 56
    elif text_level == 2:  # subtitle
        ratio = 0.60
        min_pt, max_pt = 12, 40
    else:  # body
        ratio = 0.52
        # 8-10pt is often unreadable for screenshots; raise baseline.
        min_pt, max_pt = 11, 30

    pt = max(min_pt, min(max_pt, box_h_pt * ratio))

    # Shrink by estimated wrapping
    # Rough width model: latin ~0.5-0.6em, CJK closer to 1.0em; use 0.75 to be conservative.
    char_count = len(text)
    if char_count > 0:
        chars_per_line = max(1.0, box_w_pt / (pt * 0.75))
        lines_needed = max(1, int(math.ceil(char_count / chars_per_line)))

        max_lines = max(1.0, box_h_pt / (pt * max(0.8, float(line_spacing))))
        if lines_needed > max_lines:
            pt = pt * (max_lines / lines_needed)

    pt = max(min_pt, min(max_pt, pt))
    return int(max(6, round(pt)))


def _normalize_text_bbox_px(
    bbox_px: List[int],
    *,
    canvas_h_px: int,
    text: str,
    text_level: int,
) -> List[int]:
    """
    MinerU sometimes returns extremely thin text bboxes (few pixels tall),
    which makes any font-size estimation unusable (text becomes invisible).

    This function inflates bbox height to a minimum readable band while keeping
    the horizontal placement. It is a pragmatic fix for OCR/layout bbox noise.
    
    UPDATE: Now intelligently calculates required height for long text based on 
    estimated line wrapping, preventing font from being crushed to 8pt.

    Args:
        bbox_px: [x1, y1, x2, y2] in slide-canvas px
        canvas_h_px: slide canvas height in px (target_h_px for FIGURE, slide_h_px for PDF)
        text_level: 1(title) / 2(subtitle) / 3(body)
    """
    x1, y1, x2, y2 = bbox_px
    h = max(1, y2 - y1)
    w = max(1, x2 - x1)

    # Base minimum bbox height in px on slide canvas
    if text_level == 1:
        min_h = 30
    elif text_level == 2:
        min_h = 22
    else:
        min_h = 16 if len((text or "").strip()) >= 20 else 12

    # For long text, dynamically calculate required height based on wrapping
    text_len = len((text or "").strip())
    if text_len > 50 and text_level == 3:  # Long body text
        # Target font size for estimation: 10pt (~13.3px at 96dpi)
        target_pt = 10
        target_px = target_pt * 96.0 / 72.0
        
        # Estimate chars per line (conservative: avg char width = 0.75 * font_size)
        chars_per_line = max(1.0, w / (target_px * 0.75))
        
        # Estimate lines needed
        lines_needed = max(1, int(math.ceil(text_len / chars_per_line)))
        
        # Calculate required height (lines × line_height, line_height = font_size × 1.2)
        line_height_px = target_px * 1.2
        required_h = int(lines_needed * line_height_px)
        
        # Use the larger of base min_h and calculated required_h
        min_h = max(min_h, required_h)
        
        # Cap at reasonable max to avoid excessive boxes
        min_h = min(min_h, int(canvas_h_px * 0.4))
    
    # Log for diagnosis
    log.info(f"[pdf2ppt_opt] Text BBox Normalization: text='{text[:50]}...', level={text_level}, original_bbox={bbox_px}, original_h={h}, min_h={min_h}")

    if h >= min_h:
        return bbox_px

    cy = (y1 + y2) / 2.0
    new_y1 = int(round(cy - min_h / 2.0))
    new_y2 = int(round(cy + min_h / 2.0))

    # Clamp to canvas
    new_y1 = max(0, new_y1)
    new_y2 = min(int(canvas_h_px), new_y2)
    if new_y2 <= new_y1:
        new_y2 = min(int(canvas_h_px), new_y1 + 1)

    # Log normalized bbox
    log.info(f"[pdf2ppt_opt] Normalized BBox: new_bbox=[{x1}, {new_y1}, {x2}, {new_y2}], new_h={new_y2 - new_y1}")

    return [x1, new_y1, x2, new_y2]


def _calculate_font_size(text: str, bbox: List[int], text_level: int = None) -> int:
    """
    Stable heuristic estimation (default).

    Previously we used `utils.calculate_font_size()` which directly maps bbox pixel height to pt.
    That is unstable across different slide scaling strategies (PDF vs FIGURE).
    Now we estimate using slide physical geometry (inches -> pt) so it is consistent.
    """
    lvl = text_level or 3
    return int(_estimate_font_size_from_slide_geometry_pt(text=text, bbox_px=bbox, text_level=lvl))

def _add_smart_textbox(
    slide,
    text: str,
    bbox: List[int],
    text_level: int,
    color: RGBColor = None,
    *,
    slide_w_px: int = None,
    slide_h_px: int = None,
    enable_render_fit: bool = False,
    fitter_dpi: int = 96,
):
    """
    Add a text box with stable wrapping + optional render-fit sizing (LibreOffice).
    """
    left_in = pixels_to_inches(bbox[0])
    top_in = pixels_to_inches(bbox[1])
    width_in = pixels_to_inches(bbox[2] - bbox[0])
    height_in = pixels_to_inches(bbox[3] - bbox[1])

    # Baseline heuristic (stable across PDF/FIGURE)
    font_size = _calculate_font_size(text, bbox, text_level)

    # Render-fit sizing (LibreOffice) - best for titles/short text.
    # IMPORTANT: upper_pt must be in PT, not PX.
    if enable_render_fit and slide_w_px and slide_h_px:
        try:
            style = TextFitStyle(
                font_name="Arial",
                bold=bool(text_level == 1),
                line_spacing=1.0,
                margin_px=2,  # keep a tiny margin to avoid borderline overflow
            )

            box_h_px = max(1, bbox[3] - bbox[1])
            box_h_in = pixels_to_inches(box_h_px, dpi=fitter_dpi)
            upper_pt = max(8, int(box_h_in * 72.0 * 0.95))

            font_size = DEFAULT_FITTER.fit_font_size_pt(
                text=text,
                bbox_px=(bbox[0], bbox[1], bbox[2], bbox[3]),
                slide_w_px=int(slide_w_px),
                slide_h_px=int(slide_h_px),
                style=style,
                lower_pt=8 if text_level != 1 else 14, # Allow body text to go down to 8pt
                upper_pt=upper_pt,
                tolerance_px=2,
                max_iter=15,
            )
            log.info(f"[pdf2ppt_opt] Render Fit: text='{text[:50]}...', bbox={bbox}, calculated_font_size={font_size}")
        except Exception as e:
            log.warning(f"[pdf2ppt_opt] render-fit failed, fallback to heuristic: {e}")

    # Create textbox
    textbox = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    )

    tf = textbox.text_frame
    tf.word_wrap = True

    # Titles: keep stable layout. Body: allow auto-shrink to prevent overflow.
    # UPDATE: Since we use strict render-fitting now, we MUST disable auto-fit 
    # to prevent PPT from overriding our calculated font size (e.g. forcing 8pt).
    tf.auto_size = MSO_AUTO_SIZE.NONE

    # Keep margins in sync with fitter margin_px=2 (dpi=96 => ~0.0208in)
    m_in = pixels_to_inches(2, dpi=fitter_dpi)
    tf.margin_left = Inches(m_in)
    tf.margin_right = Inches(m_in)
    tf.margin_top = Inches(m_in)
    tf.margin_bottom = Inches(m_in)

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(int(font_size))
    try:
        p.line_spacing = 1.0
    except Exception:
        pass
    
    # Font style
    # Use a standard font or keep default
    # p.font.name = "Arial" 
    
    if color:
        p.font.color.rgb = color
    else:
        p.font.color.rgb = RGBColor(0, 0, 0)
        
    if text_level == 1:
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
    elif text_level == 2:
        p.font.bold = True

@register("pdf2ppt_optimized")
def create_pdf2ppt_optimized_graph() -> GenericGraphBuilder:
    builder = GenericGraphBuilder(state_model=Paper2FigureState, entry_point="_start_")

    def _init_result_path(state: Paper2FigureState) -> Paper2FigureState:
        _ensure_result_path(state)
        return state

    async def pdf_to_images_node(state: Paper2FigureState) -> Paper2FigureState:
        if state.request.input_type == "FIGURE":
            img_path = state.request.input_content
            # 强制开启 AI 编辑，以便在转 PPT 过程中去除背景文字
            state.use_ai_edit = True
            if img_path and os.path.exists(img_path):
                state.slide_images = [img_path]
            return state

        pdf_path = getattr(state, "pdf_file", None)
        if not pdf_path:
            log.error("[pdf2ppt_opt] state.pdf_file is empty")
            return state

        base_dir = Path(_ensure_result_path(state))
        img_dir = base_dir / "slides_png"
        image_paths = ppt_tool.pdf_to_images(pdf_path, str(img_dir))
        state.slide_images = image_paths
        return state

    # --- Reused Nodes Logic from parallel wf ---
    
    async def slides_mineru_node(state: Paper2FigureState) -> Paper2FigureState:
        image_paths: List[str] = getattr(state, "slide_images", []) or []
        base_dir = Path(_ensure_result_path(state))
        mineru_dir = base_dir / "mineru_pages"
        mineru_dir.mkdir(parents=True, exist_ok=True)
        port = getattr(getattr(state, "request", None), "mineru_port", 8010)

        mineru_pages: List[Dict[str, Any]] = []

        for page_idx, img_path in enumerate(image_paths):
            try:
                out_dir = mineru_dir / f"page_{page_idx+1:03d}"
                out_dir.mkdir(parents=True, exist_ok=True)
                
                mineru_items = await recursive_mineru_layout(
                    image_path=str(img_path),
                    port=port,
                    max_depth=3,
                    output_dir=str(out_dir),
                )
                mineru_pages.append({
                    "page_idx": page_idx,
                    "blocks": mineru_items,
                    "path": img_path,
                    "mineru_output_dir": str(out_dir)
                })
            except Exception as e:
                log.error(f"[pdf2ppt_opt][MinerU] page#{page_idx+1} failed: {e}")
                mineru_pages.append({
                    "page_idx": page_idx,
                    "blocks": [],
                    "path": img_path,
                })

        state.mineru_pages = mineru_pages
        return state

    async def slides_sam_node(state: Paper2FigureState) -> Paper2FigureState:
        image_paths: List[str] = getattr(state, "slide_images", []) or []
        base_dir = _ensure_result_path(state)
        sam_pages = await asyncio.to_thread(_run_sam_on_pages, image_paths, base_dir)
        state.sam_pages = sam_pages
        return state

    async def slides_layout_bg_remove_node(state: Paper2FigureState, sam_pages: List[Dict[str, Any]] = None) -> Paper2FigureState:
        if sam_pages is None:
            sam_pages = getattr(state, "sam_pages", []) or []
        
        base_dir = Path(_ensure_result_path(state))
        icons_dir = base_dir / "sam_icons"
        icons_dir.mkdir(parents=True, exist_ok=True)
        model_path = getattr(getattr(state, "request", None), "bg_rm_model", None)

        def _sync_bg_remove():
            processed = 0
            for p in sam_pages:
                page_idx = p.get("page_idx", 0)
                for it in p.get("layout_items", []):
                    png_path = it.get("png_path")
                    if not png_path or not os.path.exists(png_path): continue
                    
                    try:
                        original_stem = Path(png_path).stem
                        output_filename = f"page_{page_idx+1:03d}_{original_stem}_bg_removed.png"
                        output_path = icons_dir / output_filename
                        
                        req = {"image_path": png_path, "output_dir": str(icons_dir)}
                        if model_path: req["model_path"] = model_path
                        
                        fg_path = local_tool_for_bg_remove(req)
                        
                        if fg_path and os.path.exists(fg_path):
                            fg_path_obj = Path(fg_path)
                            if fg_path_obj.name != output_filename:
                                new_fg_path = fg_path_obj.parent / output_filename
                                fg_path_obj.rename(new_fg_path)
                                fg_path = str(new_fg_path)
                            it["fg_png_path"] = fg_path
                        else:
                            it["fg_png_path"] = png_path
                        processed += 1
                    except Exception as e:
                        log.error(f"[pdf2ppt_opt][bg_rm] failed for {png_path}: {e}")
                        it["fg_png_path"] = png_path

            try:
                if model_path: free_bg_rm_model(model_path=model_path)
            except Exception: pass
            return processed

        await asyncio.to_thread(_sync_bg_remove)
        state.sam_pages = sam_pages
        return state

    async def parallel_processing_node(state: Paper2FigureState) -> Paper2FigureState:
        import copy
        import time
        
        start_time = time.time()
        
        async def mineru_branch():
            branch_state = copy.copy(state)
            result = await slides_mineru_node(branch_state)
            return ("mineru", result)
        
        async def sam_branch():
            branch_state = copy.copy(state)
            branch_state = await slides_sam_node(branch_state)
            sam_pages = getattr(branch_state, "sam_pages", [])
            branch_state = await slides_layout_bg_remove_node(branch_state, sam_pages=sam_pages)
            return ("sam", branch_state)
        
        results = await asyncio.gather(mineru_branch(), sam_branch(), return_exceptions=True)
        
        for r in results:
            if isinstance(r, Exception):
                log.error(f"[pdf2ppt_opt] Branch failed: {r}")
                continue
            
            branch_name, branch_state = r
            if branch_name == "mineru":
                state.mineru_pages = getattr(branch_state, "mineru_pages", None)
            elif branch_name == "sam":
                state.sam_pages = getattr(branch_state, "sam_pages", None)
        
        log.info(f"[pdf2ppt_opt] Parallel processing done in {time.time() - start_time:.2f}s")
        return state

    async def slides_ppt_generation_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        Final PPT Generation using Hybrid Logic + Dynamic Sizing + Smart Style
        """
        # Helper: API Retry
        async def _call_image_api_with_retry(coro_factory, retries: int = 3, delay: float = 1.0) -> bool:
            last_err = None
            for attempt in range(1, retries + 1):
                try:
                    await coro_factory()
                    return True
                except Exception as e:
                    last_err = e
                    log.error(f"[pdf2ppt_opt] image api failed attempt {attempt}/{retries}: {e}")
                    await asyncio.sleep(delay)
            log.error(f"[pdf2ppt_opt] image api failed after {retries} attempts: {last_err}")
            return False

        sam_pages = getattr(state, "sam_pages", []) or []
        mineru_pages = getattr(state, "mineru_pages", []) or []
        
        if not mineru_pages:
            log.error("[pdf2ppt_opt] No MinerU pages found! Aborting.")
            return state

        # Create dict indices
        sam_dict = {p.get("page_idx"): p for p in sam_pages}
        
        # Determine PPT size from first image
        first_img_path = getattr(state, "slide_images", [])[0] if getattr(state, "slide_images", []) else None

        prs = Presentation()

        # ================= 固定 PPT 物理尺寸 & 统一画布分辨率 =================
        # 避免超大输入图导致 PPT 宽度 > 56 英寸 (pptx 限制)，同时让字号/布局有稳定参照系。
        # FIGURE / PDF 都统一到固定画布，再映射 bbox。
        is_figure_mode = bool(getattr(getattr(state, "request", None), "input_type", None) == "FIGURE")

        # 物理尺寸：16:9 宽屏，限制在常见 PPT 范围内（不超过 56 in）
        target_w_in, target_h_in = 13.333, 7.5  # 16:9 widescreen
        
        # 画布像素分辨率：统一使用 1280x720 (13.333*96, 7.5*96)，方便后续 bbox 映射和字号估算
        # 如果使用 1920x1080，pixels_to_inches 会算出 20英寸 宽，导致内容溢出画布
        canvas_w_px = 1280
        canvas_h_px = 720
        
        target_w_px = canvas_w_px
        target_h_px = canvas_h_px

        slide_w_px = canvas_w_px
        slide_h_px = canvas_h_px

        prs.slide_width = Inches(target_w_in)
        prs.slide_height = Inches(target_h_in)
        log.info(
            f"[pdf2ppt_opt] PPT Size fixed to 16:9: {canvas_w_px}x{canvas_h_px} px (virtual), "
            f"{target_w_in}x{target_h_in} in; is_figure_mode={is_figure_mode}"
        )

        # FIGURE 模式下，后面 _map_bbox_px 会把原图 contain 到该画布；
        # 非 FIGURE 模式，我们会拉伸映射 (Stretch) 以匹配背景的填充方式。

        base_dir = Path(_ensure_result_path(state))

        # --- Pre-Phase: AI Background Generation ---
        use_ai_bg = bool(getattr(state, "use_ai_edit", False))
        page_bg_map = {} # page_idx -> bg_path
        ai_coroutines = []
        
        if use_ai_bg:
            log.info(f"[pdf2ppt_opt] AI Edit Enabled. Preparing background cleaning tasks...")
            masks_dir = base_dir / "masks"
            masks_dir.mkdir(parents=True, exist_ok=True)
            bg_dir = base_dir / "clean_backgrounds"
            bg_dir.mkdir(parents=True, exist_ok=True)
            
            # API Config
            req_cfg = getattr(state, "request", None) or {}
            if not isinstance(req_cfg, dict):
                req_cfg = req_cfg.__dict__ if hasattr(req_cfg, "__dict__") else {}
            api_key = req_cfg.get("api_key") or os.getenv("DF_API_KEY")
            api_url = req_cfg.get("chat_api_url") or "https://api.apiyi.com"
            model_name = req_cfg.get("gen_fig_model") or "gemini-3-pro-image-preview"
            
            if api_key:
                for page_data in mineru_pages:
                    page_idx = page_data.get("page_idx", 0)
                    img_path = page_data.get("path")
                    if not img_path or not os.path.exists(img_path):
                        continue

                    try:
                        clean_bg_path = bg_dir / f"clean_bg_{page_idx+1:03d}.png"
                        page_bg_map[page_idx] = str(clean_bg_path)

                        # New Prompt: Single Image Edit
                        prompt = "Remove all text from the image, keeping only the background, figures, and icons. Do not change the layout or style. 去除文字，只保留底色 图像 图标"

                        async def _run_ai(ip=img_path, op=str(clean_bg_path)):
                            await _call_image_api_with_retry(
                                lambda: generate_or_edit_and_save_image_async(
                                    prompt=prompt,
                                    image_path=ip,
                                    save_path=op,
                                    api_url=api_url,
                                    api_key=api_key,
                                    model=model_name,
                                    use_edit=True,
                                    resolution="2K",
                                    timeout=300,
                                )
                            )

                        ai_coroutines.append(_run_ai())
                    except Exception as e:
                        log.error(f"[pdf2ppt_opt] AI Edit setup failed page#{page_idx}: {e}")

                if ai_coroutines:
                    log.info(f"[pdf2ppt_opt] Waiting for {len(ai_coroutines)} AI tasks...")
                    await asyncio.gather(*ai_coroutines, return_exceptions=True)
            else:
                log.warning("[pdf2ppt_opt] use_ai_edit is True but no API Key found.")

        # Coordinate mapping:
        # - non-FIGURE (PDF): Stretch mapping (匹配背景的拉伸填充)
        # - FIGURE: scale original image into standard 16:9 canvas (contain) with centering
        def _map_bbox_px(b: List[int], src_w_px: int, src_h_px: int) -> List[int]:
            x1, y1, x2, y2 = b
            
            if src_w_px <= 0 or src_h_px <= 0:
                return b

            if not is_figure_mode:
                # PDF Mode: Stretch to fill
                # 背景图是用 prs.slide_width/height 拉伸铺满的，所以前景也要同样拉伸
                sx = target_w_px / src_w_px
                sy = target_h_px / src_h_px
                return [
                    int(round(x1 * sx)),
                    int(round(y1 * sy)),
                    int(round(x2 * sx)),
                    int(round(y2 * sy)),
                ]
            else:
                # FIGURE Mode: Contain (Fit within canvas)
                s = min(target_w_px / src_w_px, target_h_px / src_h_px)
                dx = (target_w_px - src_w_px * s) / 2.0
                dy = (target_h_px - src_h_px * s) / 2.0

                return [
                    int(round(x1 * s + dx)),
                    int(round(y1 * s + dy)),
                    int(round(x2 * s + dx)),
                    int(round(y2 * s + dy)),
                ]

        for page_data in mineru_pages:
            page_idx = page_data.get("page_idx", 0)
            mineru_blocks = page_data.get("blocks", [])
            img_path = page_data.get("path")
            
            # Get corresponding SAM data
            sam_data = sam_dict.get(page_idx, {})
            sam_items = sam_data.get("layout_items", [])

            slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank

            # Per-page source image size (used for mapping)
            src_w_px, src_h_px = slide_w_px, slide_h_px
            try:
                with Image.open(img_path) as _img:
                    src_w_px, src_h_px = _img.size
            except Exception:
                pass

            # --- 1. Background ---
            clean_bg = page_bg_map.get(page_idx)
            if clean_bg and os.path.exists(clean_bg):
                try:
                    # In FIGURE mode we also want the background to be scaled+centered into canvas,
                    # to keep alignment with mapped bboxes.
                    if is_figure_mode:
                        s_bg = min(target_w_px / src_w_px, target_h_px / src_h_px)
                        w_bg = int(round(src_w_px * s_bg))
                        h_bg = int(round(src_h_px * s_bg))
                        dx_bg = int(round((target_w_px - w_bg) / 2.0))
                        dy_bg = int(round((target_h_px - h_bg) / 2.0))
                        slide.shapes.add_picture(
                            clean_bg,
                            Inches(pixels_to_inches(dx_bg)),
                            Inches(pixels_to_inches(dy_bg)),
                            Inches(pixels_to_inches(w_bg)),
                            Inches(pixels_to_inches(h_bg)),
                        )
                    else:
                        slide.shapes.add_picture(clean_bg, 0, 0, prs.slide_width, prs.slide_height)
                except Exception as e:
                    log.error(f"[pdf2ppt_opt] Failed to add clean bg: {e}")
                    # Fallback to white
                    bg = slide.background
                    fill = bg.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(255, 255, 255)
            else:
                # White background (Reconstruction Mode)
                bg = slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(255, 255, 255)

            # --- 2. Filter SAM Items (Icons/Logos) ---
            # We filter out SAM items that overlap heavily with MinerU text/image blocks
            final_sam_items = []
            
            # Convert MinerU blocks to pixel bboxes for overlap check (mapped to canvas when FIGURE)
            mineru_bboxes_px = []
            for blk in mineru_blocks:
                bbox_n = blk.get("bbox")
                if bbox_n:
                    x1 = int(bbox_n[0] * src_w_px)
                    y1 = int(bbox_n[1] * src_h_px)
                    x2 = int(bbox_n[2] * src_w_px)
                    y2 = int(bbox_n[3] * src_h_px)
                    mineru_bboxes_px.append(_map_bbox_px([x1, y1, x2, y2], src_w_px, src_h_px))

            def _get_iou(boxA, boxB):
                xA = max(boxA[0], boxB[0])
                yA = max(boxA[1], boxB[1])
                xB = min(boxA[2], boxB[2])
                yB = min(boxA[3], boxB[3])
                interArea = max(0, xB - xA) * max(0, yB - yA)
                boxAArea = (boxA[2] - boxA[0]) * (boxA[3] - boxA[1])
                if boxAArea == 0: return 0
                return interArea / boxAArea

            for item in sam_items:
                s_bbox_raw = item.get("bbox_px")
                if not s_bbox_raw: continue

                # FIX: Map SAM bbox (source coords) to Canvas coords for correct IOU check
                # MinerU bboxes are already mapped to canvas in mineru_bboxes_px
                s_bbox_mapped = _map_bbox_px(
                    [s_bbox_raw[0], s_bbox_raw[1], s_bbox_raw[2], s_bbox_raw[3]], 
                    src_w_px, 
                    src_h_px
                )
                
                # Check overlap with MinerU blocks
                is_overlap = False
                for m_bbox in mineru_bboxes_px:
                    if _get_iou(s_bbox_mapped, m_bbox) > 0.5:
                        is_overlap = True
                        break
                
                if not is_overlap:
                    final_sam_items.append(item)

            # Render SAM items (mapped in FIGURE mode)
            for item in final_sam_items:
                fg_path = item.get("fg_png_path")
                raw_path = item.get("png_path")
                
                # Log for diagnosis
                log.info(f"[pdf2ppt_opt] SAM Item: raw_path={raw_path}, fg_path={fg_path}, exists_fg={os.path.exists(fg_path) if fg_path else False}")
                
                # Only render if fg_path exists (background removed successfully)
                if fg_path and os.path.exists(fg_path):
                    ipath = fg_path
                else:
                    log.warning(f"[pdf2ppt_opt] Skipping SAM item due to missing fg_png_path: {raw_path}")
                    continue
                
                bbox = item.get("bbox_px")
                if not bbox:
                    continue
                mb = _map_bbox_px([bbox[0], bbox[1], bbox[2], bbox[3]], src_w_px, src_h_px)
                slide.shapes.add_picture(
                    ipath,
                    Inches(pixels_to_inches(mb[0])),
                    Inches(pixels_to_inches(mb[1])),
                    Inches(pixels_to_inches(mb[2] - mb[0])),
                    Inches(pixels_to_inches(mb[3] - mb[1])),
                )

            # --- 3. Process MinerU Blocks (Text & Images) ---
            for blk in mineru_blocks:
                btype = (blk.get("type") or "").lower()
                bbox_n = blk.get("bbox")
                if not bbox_n: continue
                
                # Convert to pixels (source image space) then map to slide canvas if FIGURE
                x1 = int(bbox_n[0] * src_w_px)
                y1 = int(bbox_n[1] * src_h_px)
                x2 = int(bbox_n[2] * src_w_px)
                y2 = int(bbox_n[3] * src_h_px)
                bbox_px = _map_bbox_px([x1, y1, x2, y2], src_w_px, src_h_px)
                
                if x2 <= x1 or y2 <= y1: continue

                # A. Images/Tables -> Render directly (using MinerU crops or fallback)
                if btype in ['image', 'figure', 'table', 'formula']:
                    # Try to find existing image path
                    ipath = blk.get("img_path")
                    # If not in block, maybe MinerU saved sub-images. 
                    # For now, simple fallback: crop from original
                    if not ipath or not os.path.exists(ipath):
                        # Use fallback crop
                        fallback_dir = base_dir / "mineru_fallback"
                        fallback_dir.mkdir(parents=True, exist_ok=True)
                        ipath = str(fallback_dir / f"crop_{page_idx}_{id(blk)}.png")
                        if not os.path.exists(ipath):
                            try:
                                with Image.open(img_path) as page_img:
                                    crop = page_img.crop((x1, y1, x2, y2))
                                    crop.save(ipath)
                            except:
                                ipath = None
                    
                    if ipath and os.path.exists(ipath):
                         slide.shapes.add_picture(
                            ipath,
                            Inches(pixels_to_inches(bbox_px[0])),
                            Inches(pixels_to_inches(bbox_px[1])),
                            Inches(pixels_to_inches(bbox_px[2] - bbox_px[0])),
                            Inches(pixels_to_inches(bbox_px[3] - bbox_px[1])),
                        )

                # B. Text/Title -> Use MinerU Text
                elif btype in ['text', 'title', 'header', 'footer', 'reference', 'list']:
                    text_level = 1 if btype in ['title', 'header'] else 3
                    if btype == 'text':
                        text_level = 3
                    
                    final_text = blk.get("text") or blk.get("content") or ""
                    
                    if not final_text.strip():
                        continue

                    # canvas size in px for bbox normalization & fitter coordinate system
                    canvas_w_px_for_fit = target_w_px if is_figure_mode else slide_w_px
                    canvas_h_px_for_fit = target_h_px if is_figure_mode else slide_h_px

                    # 1. Fix pathological thin bboxes -> invisible text
                    bbox_px = _normalize_text_bbox_px(
                        bbox_px,
                        canvas_h_px=int(canvas_h_px_for_fit),
                        text=final_text,
                        text_level=text_level,
                    )

                    # 2. Color Sampling
                    text_color = _get_dominant_color(img_path, bbox_px)
                    
                    # 3. Add Text Box
                    # Enable render-fit for ALL text to ensure it fits within bbox
                    enable_render_fit = True
                    _add_smart_textbox(
                        slide,
                        final_text,
                        bbox_px,
                        text_level,
                        text_color,
                        slide_w_px=canvas_w_px_for_fit,
                        slide_h_px=canvas_h_px_for_fit,
                        enable_render_fit=enable_render_fit,
                        fitter_dpi=96,
                    )

        # Save PPT
        ppt_path = base_dir / "pdf2ppt_optimized_output.pptx"
        prs.save(str(ppt_path))
        state.ppt_path = str(ppt_path)
        log.info(f"[pdf2ppt_opt] PPT generated at: {ppt_path}")
        
        return state

    nodes = {
        "_start_": _init_result_path,
        "pdf_to_images": pdf_to_images_node,
        "parallel_processing": parallel_processing_node,
        "slides_ppt_generation": slides_ppt_generation_node,
        "_end_": lambda state: state,
    }

    edges = [
        ("pdf_to_images", "parallel_processing"),
        ("parallel_processing", "slides_ppt_generation"),
        ("slides_ppt_generation", "_end_"),
    ]

    builder.add_nodes(nodes).add_edges(edges)
    builder.add_edge("_start_", "pdf_to_images")
    return builder
