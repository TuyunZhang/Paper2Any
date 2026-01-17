from __future__ import annotations

import hashlib
import os
import tempfile
import subprocess
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, Optional, Tuple

import fitz  # PyMuPDF
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE

from dataflow_agent.logger import get_logger
from dataflow_agent.utils import pixels_to_inches

log = get_logger(__name__)


@dataclass(frozen=True)
class TextFitStyle:
    font_name: str = "Arial"
    bold: bool = False
    line_spacing: float = 1.0
    margin_px: int = 0  # internal margin in px (converted to inches using dpi)


class PptTextFitter:
    """
    Fit text size by real rendering:
    1) Create a minimal pptx with one textbox at the target bbox
    2) Use LibreOffice (soffice) to convert pptx -> pdf
    3) Render pdf to image via PyMuPDF
    4) Detect whether text pixels exceed bbox (with small tolerance)
    5) Binary search the largest font size that fits

    Notes:
    - This is slower than heuristic estimation; use it selectively (e.g. titles).
    - It works even when PPT's own wrapping differs from our heuristics, because it measures output.
    """

    def __init__(self, dpi: int = 96):
        self.dpi = dpi
        self._cache: Dict[str, int] = {}

    def fit_font_size_pt(
        self,
        *,
        text: str,
        bbox_px: Tuple[int, int, int, int],
        slide_w_px: int,
        slide_h_px: int,
        style: TextFitStyle,
        lower_pt: int = 6,
        upper_pt: Optional[int] = None,
        tolerance_px: int = 2,
        max_iter: int = 10,
    ) -> int:
        text = (text or "").strip()
        if not text:
            return 12

        x1, y1, x2, y2 = bbox_px
        box_w = max(1, x2 - x1)
        box_h = max(1, y2 - y1)

        # Default upper bound from height; titles can be huge if box is tall.
        if upper_pt is None:
            upper_pt = max(lower_pt + 1, int(box_h * 0.8))

        cache_key = self._make_cache_key(
            text=text,
            box_w=box_w,
            box_h=box_h,
            style=style,
            upper_pt=upper_pt,
            lower_pt=lower_pt,
        )
        cached = self._cache.get(cache_key)
        if cached is not None:
            return cached

        lo, hi = int(lower_pt), int(upper_pt)
        best = lo

        for _ in range(max_iter):
            if lo > hi:
                break
            mid = (lo + hi) // 2
            ok = self._render_and_check_fit(
                text=text,
                bbox_px=bbox_px,
                slide_w_px=slide_w_px,
                slide_h_px=slide_h_px,
                style=style,
                font_pt=mid,
                tolerance_px=tolerance_px,
            )
            if ok:
                best = mid
                lo = mid + 1
            else:
                hi = mid - 1

        self._cache[cache_key] = best
        return best

    def _make_cache_key(
        self,
        *,
        text: str,
        box_w: int,
        box_h: int,
        style: TextFitStyle,
        upper_pt: int,
        lower_pt: int,
    ) -> str:
        h = hashlib.sha1()
        h.update(text.encode("utf-8", errors="ignore"))
        h.update(f"|{box_w}x{box_h}|{style.font_name}|{int(style.bold)}|{style.line_spacing}|{style.margin_px}|{lower_pt}-{upper_pt}".encode())
        return h.hexdigest()

    def _render_and_check_fit(
        self,
        *,
        text: str,
        bbox_px: Tuple[int, int, int, int],
        slide_w_px: int,
        slide_h_px: int,
        style: TextFitStyle,
        font_pt: int,
        tolerance_px: int,
    ) -> bool:
        with tempfile.TemporaryDirectory(prefix="ppt_text_fit_") as td:
            tmpdir = Path(td)
            pptx_path = tmpdir / "fit.pptx"
            pdf_path = tmpdir / "fit.pdf"

            self._create_minimal_pptx(
                pptx_path=pptx_path,
                text=text,
                bbox_px=bbox_px,
                slide_w_px=slide_w_px,
                slide_h_px=slide_h_px,
                style=style,
                font_pt=font_pt,
            )

            ok = self._convert_pptx_to_pdf(pptx_path=pptx_path, out_dir=tmpdir)
            if not ok or not pdf_path.exists():
                # If conversion fails, be conservative: treat as not fit.
                return False

            img = self._render_pdf_first_page(pdf_path=pdf_path)
            if img is None:
                return False

            return self._check_text_pixels_within_bbox(
                rendered_rgb=img,
                bbox_px=bbox_px,
                tolerance_px=tolerance_px,
            )

    def _create_minimal_pptx(
        self,
        *,
        pptx_path: Path,
        text: str,
        bbox_px: Tuple[int, int, int, int],
        slide_w_px: int,
        slide_h_px: int,
        style: TextFitStyle,
        font_pt: int,
    ) -> None:
        prs = Presentation()
        prs.slide_width = Inches(pixels_to_inches(slide_w_px, dpi=self.dpi))
        prs.slide_height = Inches(pixels_to_inches(slide_h_px, dpi=self.dpi))

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        x1, y1, x2, y2 = bbox_px
        left = Inches(pixels_to_inches(x1, dpi=self.dpi))
        top = Inches(pixels_to_inches(y1, dpi=self.dpi))
        width = Inches(pixels_to_inches(max(1, x2 - x1), dpi=self.dpi))
        height = Inches(pixels_to_inches(max(1, y2 - y1), dpi=self.dpi))

        shape = slide.shapes.add_textbox(left, top, width, height)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE

        # Internal margins -> keep in sync with workflow for stable fitting
        margin_in = pixels_to_inches(style.margin_px, dpi=self.dpi)
        tf.margin_left = Inches(margin_in)
        tf.margin_right = Inches(margin_in)
        tf.margin_top = Inches(margin_in)
        tf.margin_bottom = Inches(margin_in)

        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(int(font_pt))
        p.font.bold = bool(style.bold)
        if style.font_name:
            p.font.name = style.font_name
        try:
            p.line_spacing = float(style.line_spacing)
        except Exception:
            pass

        prs.save(str(pptx_path))

    def _convert_pptx_to_pdf(self, *, pptx_path: Path, out_dir: Path) -> bool:
        soffice = os.environ.get("SOFFICE_BIN") or "soffice"
        cmd = [
            soffice,
            "--headless",
            "--nologo",
            "--nofirststartwizard",
            "--convert-to",
            "pdf",
            "--outdir",
            str(out_dir),
            str(pptx_path),
        ]
        try:
            res = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
            if res.returncode != 0:
                log.warning(f"[ppt_text_fit] soffice convert failed: {res.stderr or res.stdout}")
                return False
            return True
        except Exception as e:
            log.warning(f"[ppt_text_fit] soffice convert exception: {e}")
            return False

    def _render_pdf_first_page(self, *, pdf_path: Path) -> Optional[np.ndarray]:
        try:
            doc = fitz.open(str(pdf_path))
            try:
                page = doc[0]
                # Render at 1x; since the ppt slide is in inches with dpi mapping, this is stable enough.
                pix = page.get_pixmap(alpha=False)
                img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, 3)
                return img
            finally:
                doc.close()
        except Exception as e:
            log.warning(f"[ppt_text_fit] render pdf failed: {e}")
            return None

    def _check_text_pixels_within_bbox(
        self,
        *,
        rendered_rgb: np.ndarray,
        bbox_px: Tuple[int, int, int, int],
        tolerance_px: int,
    ) -> bool:
        """
        Detect whether non-background pixels inside bbox spill outside bbox.
        Strategy:
        - Assume background is mostly white (as we don't add background in minimal pptx).
        - Threshold to detect ink pixels.
        - Find bounding box of ink pixels and ensure it's within bbox (+ tolerance).
        """
        x1, y1, x2, y2 = bbox_px
        H, W = rendered_rgb.shape[:2]

        # Clamp bbox to rendered image size
        x1c = max(0, min(W - 1, x1))
        y1c = max(0, min(H - 1, y1))
        x2c = max(0, min(W, x2))
        y2c = max(0, min(H, y2))
        if x2c <= x1c or y2c <= y1c:
            return True

        # Convert to grayscale
        gray = (0.299 * rendered_rgb[:, :, 0] + 0.587 * rendered_rgb[:, :, 1] + 0.114 * rendered_rgb[:, :, 2]).astype(
            np.uint8
        )

        # Ink mask: anything darker than 245 considered ink (titles are dark)
        ink = gray < 245

        # Get ink bounding box
        ys, xs = np.where(ink)
        if len(xs) == 0:
            return True  # no ink => fits

        minx = int(xs.min())
        maxx = int(xs.max())
        miny = int(ys.min())
        maxy = int(ys.max())

        # Check if ink bbox is within target bbox with tolerance
        if minx < x1c - tolerance_px:
            return False
        if miny < y1c - tolerance_px:
            return False
        if maxx > x2c + tolerance_px:
            return False
        if maxy > y2c + tolerance_px:
            return False

        return True


DEFAULT_FITTER = PptTextFitter()
