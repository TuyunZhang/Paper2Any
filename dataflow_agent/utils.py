import ast
from json import JSONDecodeError, JSONDecoder
import json
import re
from typing import Any, Dict, Union, List
from pathlib import Path
from dataflow_agent.logger import get_logger
log = get_logger(__name__)

import ast
from json import JSONDecodeError, JSONDecoder
import json
import re
from typing import Any, Dict, Union, List
from pathlib import Path
import asyncio
import logging
from typing import List, Dict, Any, Optional


from math import ceil
import time
import os
import math
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

import fitz  # PyMuPDF

from PIL import Image
import pdfplumber
import uuid

from dataflow_agent.toolkits.multimodaltool.mineru_tool import run_aio_batch_two_step_extract, run_aio_two_step_extract

def get_project_root() -> Path:
    return Path(__file__).resolve().parent.parent

def robust_parse_json(
    text: str,
    *,
    merge_dicts: bool = False,
    strip_double_braces: bool = False
) -> Union[Dict[str, Any], List[Any]]:
    """
    尽量从 LLM / 日志 / jsonl / Markdown 片段中提取合法 JSON。

    参数
    ----
    text : str
        输入原始文本
    merge_dicts : bool, default False
        提取到多个对象且全部是 dict 时，是否用 dict.update 合并返回
    strip_double_braces : bool, default False
        把 '{{' / '}}' 替换成 '{' / '}'（某些模板语言会加双层花括号）

    返回
    ----
    Dict / List / List[Dict | List]
    """
    s = text.strip()

    # ---------- 预处理：剥去外层包裹 ----------
    s = _remove_markdown_fence(s)          # ```json ... ```
    s = _remove_outer_triple_quotes(s)     # ''' ... ''' / """ ... """
    s = _remove_leading_json_word(s)       # 开头一个 json/JSON 标记

    if strip_double_braces:
        s = s.replace("{{", "{").replace("}}", "}")

    # ---------- 清理注释 & 尾逗号 ----------
    s = _strip_json_comments(s)

    # ---------- 新增：清理非法控制字符 ----------
    # 移除所有 JSON 规范不允许的 ASCII 控制字符。
    # 合法的 \n, \r, \t, 和 \f, \b, \" 都不会被移除，但这里只针对不可打印的控制码。
    s = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', s)

    # ---------- 新增：转义未转义的反斜杠（修复 LaTeX 公式等问题）----------
    # 这会将所有单个反斜杠转换为双反斜杠，但保留已经正确转义的序列
    # 先保护已经转义的序列（如 \\n, \\t, \\", \\\\）
    s = s.replace('\\\\', '\x00DOUBLE_BACKSLASH\x00')  # 临时标记
    s = s.replace('\\n', '\x00NEWLINE\x00')
    s = s.replace('\\r', '\x00RETURN\x00')
    s = s.replace('\\t', '\x00TAB\x00')
    s = s.replace('\\"', '\x00QUOTE\x00')
    s = s.replace('\\/', '\x00SLASH\x00')
    s = s.replace('\\b', '\x00BACKSPACE\x00')
    s = s.replace('\\f', '\x00FORMFEED\x00')
    
    # 现在转义所有剩余的单个反斜杠
    s = s.replace('\\', '\\\\')
    
    # 恢复之前保护的序列
    s = s.replace('\x00DOUBLE_BACKSLASH\x00', '\\\\')
    s = s.replace('\x00NEWLINE\x00', '\\n')
    s = s.replace('\x00RETURN\x00', '\\r')
    s = s.replace('\x00TAB\x00', '\\t')
    s = s.replace('\x00QUOTE\x00', '\\"')
    s = s.replace('\x00SLASH\x00', '\\/')
    s = s.replace('\x00BACKSPACE\x00', '\\b')
    s = s.replace('\x00FORMFEED\x00', '\\f')

    log.debug(f'清洗完之后内容是： {s}')

    # ---------- Step-1：整体解析 ----------
    # Step-1
    try:
        result = json.loads(s)
        log.info(f"整体解析成功，类型: {type(result)}")
        return result
    except JSONDecodeError as e:
        log.warning(f"整体解析失败: {e}")

    # ---------- Step-2：尝试 JSON Lines ----------
    objs = _parse_json_lines(s)
    if objs is not None:
        return _maybe_merge(objs, merge_dicts)

    # ---------- Step-3：流式提取多个对象 ----------
    objs = _extract_json_objects(s)
    log.warning(f"提取到 {len(objs)} 个对象")
    if not objs:
        raise ValueError("Unable to locate any valid JSON fragment.")

    return _maybe_merge(objs, merge_dicts)


# ======================================================================
#                            工具函数
# ======================================================================

_fence_pat = re.compile(r'```[\w-]*\s*([\s\S]*?)```', re.I)
# 只匹配外层包裹的代码块（整个内容被 ``` 包裹）
_outer_fence_pat = re.compile(r'^\s*```[\w-]*\s*([\s\S]*?)```\s*$', re.I)


def _remove_markdown_fence(src: str) -> str:
    """只提取外层包裹的 ``` … ``` 内文本；若内容不是被代码块包裹则原样返回"""
    # 只处理整个内容被代码块包裹的情况，避免提取 JSON 内嵌的代码块
    match = _outer_fence_pat.match(src)
    if match:
        return match.group(1).strip()
    return src


def _remove_outer_triple_quotes(src: str) -> str:
    if (src.startswith("'''") and src.endswith("'''")) or (
        src.startswith('"""') and src.endswith('"""')
    ):
        return src[3:-3].strip()
    return src


def _remove_leading_json_word(src: str) -> str:
    return src[4:].lstrip() if src.lower().startswith("json") else src


def _strip_json_comments(src: str) -> str:
    # /* ... */  块注释
    src = re.sub(r'/\*[\s\S]*?\*/', '', src)
    # // ...     行注释，排除 URL 中的 :// 和字符串内的 //）
    src = re.sub(r'(?<![:\"\'])//.*', '', src)
    # 尾逗号 ,}
    src = re.sub(r',\s*([}\]])', r'\1', src)
    return src.strip()


# ----------------  JSON Lines ----------------
def _parse_json_lines(src: str) -> Union[List[Any], None]:
    lines = [ln.strip() for ln in src.splitlines() if ln.strip()]
    if len(lines) <= 1:          # 只有 0/1 行就不是 jsonl
        return None

    objs: List[Any] = []
    for ln in lines:
        try:
            objs.append(json.loads(ln))
        except JSONDecodeError:
            return None  # 某行不是合法 JSON，放弃 jsonl 方案
    return objs


# ------------  多对象提取（改进版） ------------
def _extract_json_objects(src: str) -> List[Any]:
    dec = JSONDecoder()
    idx, n = 0, len(src)
    objs: List[Any] = []

    while idx < n:
        m = re.search(r'[{\[]', src[idx:])
        if not m:
            break
        idx += m.start()
        try:
            obj, end = dec.raw_decode(src, idx)
            # ========== 严格性检查 ==========
            tail = src[end:].lstrip()
            # 允许结束、逗号、换行、右括号、右中括号
            if tail and tail[0] not in ',]}>\n\r':
                idx += 1  # 可能是误判，如  {"a":1  <-- 缺 }
                continue
            objs.append(obj)
            idx = end
        except JSONDecodeError:
            idx += 1
    return objs


def _maybe_merge(objs: List[Any], merge_dicts: bool) -> Union[Any, List[Any]]:
    if len(objs) == 1:
        return objs[0]
    if merge_dicts and all(isinstance(o, dict) for o in objs):
        merged: Dict[str, Any] = {}
        for o in objs:
            merged.update(o)
        return merged
    return objs


# ========================================================================

#                           For Paper2Figure

# ========================================================================

async def run_mineru(image_path: Path, output_dir: Path) -> bool:
    """调用 mineru 并返回是否成功"""
    cmd = [
        "mineru",
        "-p", str(image_path),
        "--backend", "vlm-transformers",
        "--source", "local",
        "-o", str(output_dir)
    ]
    proc = await asyncio.create_subprocess_exec(
        *cmd,
        stdout=asyncio.subprocess.PIPE,
        stderr=asyncio.subprocess.PIPE,
    )
    stdout, stderr = await proc.communicate()

    if proc.returncode != 0:
        log.warning(f"[mineru] 运行失败: {stderr.decode(errors='ignore')}")
        return False
    
    log.info("[mineru] 执行成功")
    return True

async def replace_item_with_sub_items(
    items: List[Dict[str, Any]], 
    sub_items: List[Dict[str, Any]], 
    sub_img_path: str
) -> List[Dict[str, Any]]:
    """
    处理每个 sub_item，执行坐标变换并替换原有的 item。
    """
    # 加载图像并获取宽高
    with Image.open(sub_img_path) as sub_img:
        sub_img_width, sub_img_height = sub_img.size

    expanded_items = []  # 用来存储展开后的 sub_items
    items_to_remove = []  # 用来记录需要删除的 items 索引

    log.info(f"[replace_item_with_sub_items] 开始替换图像: {sub_img_path}")

    for i, item in enumerate(items):
        log.info(f"item[type]: {item['type']}")
        if item["type"] in ["image", "table"]:
            log.info(f"item['img_path']: {item['img_path']}")
            log.info(f"sub_img_path: {sub_img_path}") 
            if item["img_path"] == sub_img_path: # 用 sub_img_path 匹配原图的 img_path
                # 获取目标图像的 bbox
                target_bbox = item["bbox"]
                xmin, ymin, xmax, ymax = target_bbox
                # 计算比例
                width_ratio = (xmax - xmin) / sub_img_width
                height_ratio = (ymax - ymin) / sub_img_height

                log.info(f"[replace_item_with_sub_items] 找到匹配的 item，替换 bbox: {target_bbox} -> 展开为 {len(sub_items)} 个 sub_items")

                # 替换原始项中的 image 信息，展开 sub_items
                for sub_item in sub_items:
                    # 获取 sub_item 的原始 bbox 坐标
                    sub_item_bbox = sub_item["bbox"]
                    sub_xmin, sub_ymin, sub_xmax, sub_ymax = sub_item_bbox

                    # 进行坐标线性变换
                    transformed_bbox = [
                        int(sub_xmin * width_ratio + xmin),  # x 坐标变换
                        int(sub_ymin * height_ratio + ymin),  # y 坐标变换
                        int(sub_xmax * width_ratio + xmin),  # x 坐标变换
                        int(sub_ymax * height_ratio + ymin)  # y 坐标变换
                    ]

                    # 更新 sub_item 的 bbox
                    sub_item["bbox"] = transformed_bbox

                    # 将变换后的 sub_item 加入 expanded_items 列表
                    expanded_items.append(sub_item)

                # 记录需要删除的原 item 索引
                items_to_remove.append(i)

                break

    # 替换原项中的 item，删除原有的 item
    log.info(f"[replace_item_with_sub_items] 将原 item 删除, 替换为 {len(expanded_items)} 个 sub_item")

    for index in sorted(items_to_remove, reverse=True):  # 逆序删除，避免索引错乱
        del items[index]

    # 将展开的 sub_items 插入到 items 中
    items.extend(expanded_items)

    log.info(f"[replace_item_with_sub_items] 替换完成后，当前 items 长度: {len(items)}")

    return items

async def recursive_run_mineru(
    img_path: Path, 
    out_dir: Path, 
    max_depth: int = 2, 
    current_depth: int = 0
) -> List[Dict[str, Any]]:
    """递归运行 mineru，处理子图并更新fig_mask"""
    
    if current_depth > max_depth:
        return []  # 达到最大递归深度时停止递归
    
    log.info(f"[recursive_run_mineru] 当前深度 {current_depth}，处理图像: {img_path}")
    
    # 调用 mineru 处理当前图像
    ok = await run_mineru(img_path, out_dir)
    if not ok:
        return []  # 如果失败，则返回空结果

    # 寻找并读取中间结果的 JSON 文件
    content_json = locate_content_json(out_dir / img_path.stem)
    if content_json is None:
        return []  # 没有找到内容文件，返回空结果

    items = load_and_fix_items(content_json, out_dir)

    log.info(f"[recursive_run_mineru] 当前 items 长度: {len(items)}")

    # 查找子图目录并递归处理每个子图
    vlm_images_dir = out_dir / img_path.stem / 'vlm' / 'images'
    sub_images = list(vlm_images_dir.glob("*.jpg"))  # 假设子图为 .jpg 格式

    if(current_depth != max_depth):
        for sub_img_path in sub_images:  # sub_img_path 是图像路径
            log.info(f"[recursive_run_mineru] 处理子图: {sub_img_path}")

            # 获取当前 sub_img_path 对应的 sub_items
            sub_items = await recursive_run_mineru(sub_img_path, out_dir, max_depth, current_depth + 1)

            # 处理替换和坐标变换
            items = await replace_item_with_sub_items(items, sub_items, str(sub_img_path))

            log.info(f"[recursive_run_mineru] 替换完成后，当前 items 长度: {len(items)}")

    return items

def locate_content_json(output_dir: Path) -> Path | None:
    """寻找 *_middle.json 文件"""
    files = list(output_dir.rglob("*_middle.json"))
    if not files:
        log.warning(f"[mineru] 未找到 *_middle.json in {output_dir}")
        return None
    return files[0]


def load_and_fix_items(content_json: Path, output_dir: Path) -> List[Dict[str, Any]]:
    """
    读取 JSON 并提取所有文本和图片元素，修复图片路径为绝对路径
    """
    try:
        data = json.loads(content_json.read_text(encoding="utf-8"))
    except Exception as e:
        log.warning(f"[mineru] JSON 读取失败: {e}")
        return []

    # 提取基础名称
    stem = content_json.stem  # e.g. "paper1_middle"
    base_name = stem.replace("_middle", "") or stem
    
    results = []
    
    # 遍历pdf_info中的para_blocks
    if "pdf_info" in data and isinstance(data["pdf_info"], list):
        for pdf_info in data["pdf_info"]:
            if "para_blocks" in pdf_info and isinstance(pdf_info["para_blocks"], list):
                for block in pdf_info["para_blocks"]:
                    block_type = block.get("type", "")
                    
                    # 处理标题和普通文本
                    if block_type in ["title", "text", "paragraph"]:
                        # 提取文本内容
                        text_content = extract_text_from_block(block)
                        if text_content:
                            results.append({
                                "type": "text",
                                "text": text_content,
                                "bbox": block.get("bbox", []),
                                "text_level": 1 if block_type == "title" else None,
                                "page_idx": 0
                            })
                    
                    # 处理图片块
                    elif block_type in ["list", "image", "table"]:
                        # 提取图片和相关的caption
                        image_elements = extract_image_elements(block, base_name, output_dir)
                        results.extend(image_elements)
    
    return results


def extract_text_from_block(block: Dict) -> str:
    """从块中提取文本内容"""
    text_parts = []
    
    # 如果有lines字段，遍历提取
    if "lines" in block and isinstance(block["lines"], list):
        for line in block["lines"]:
            if "spans" in line and isinstance(line["spans"], list):
                for span in line["spans"]:
                    if span.get("type") == "text" and "content" in span:
                        text_parts.append(span["content"])
    
    # 如果没有lines字段，尝试直接获取content
    elif "content" in block:
        text_parts.append(block["content"])
    
    return " ".join(text_parts) if text_parts else ""


def extract_image_elements(block: Dict, base_name: str, output_dir: Path) -> List[Dict]:
    """从图片块中提取图片和相关的文本元素"""
    elements = []
    
    if "blocks" in block and isinstance(block["blocks"], list):
        for sub_block in block["blocks"]:
            sub_type = sub_block.get("type", "")
            
            # 处理图片caption
            if sub_type in ["title", "text", "paragraph","image_caption", "table_caption"]:
                caption_text = extract_text_from_block(sub_block)
                if caption_text:
                    elements.append({
                        "type": "text",
                        "text": caption_text,
                        "bbox": sub_block.get("bbox", []),
                        "text_level": None,  # caption作为普通文本
                        "page_idx": 0
                    })
            
            # 处理图片主体
            elif sub_type in ["image_body", "table_body"]:
                image_path = extract_image_path(sub_block, base_name, output_dir)
                if image_path:
                    elements.append({
                        "type": "image",
                        "img_path": str(image_path),
                        "bbox": sub_block.get("bbox", []),
                        "image_caption": [],
                        "image_footnote": [],
                        "page_idx": 0
                    })
    
    return elements


def extract_image_path(block: Dict, base_name: str, output_dir: Path) -> Optional[Path]:
    """从图片块中提取图片路径并转换为绝对路径"""
    if "lines" in block and isinstance(block["lines"], list):
        for line in block["lines"]:
            if "spans" in line and isinstance(line["spans"], list):
                for span in line["spans"]:
                    if span.get("type") in ["image", "table"] and "image_path" in span:
                        rel_path = span["image_path"]
                        if rel_path:
                            # 构建绝对路径
                            abs_path = output_dir / base_name / "vlm/images" / rel_path
                            return abs_path
    return None

def build_output_directory(image_path: Path) -> Path:
    """构造 <image_no_ext>_mineru 输出目录"""
    base = image_path.with_suffix("")
    out_dir = Path(f"{base}_mineru")
    out_dir.mkdir(parents=True, exist_ok=True)
    return out_dir


# -------------------- NEW UTILS FUNCTIONS ---------------------------------

import asyncio
from pathlib import Path
from PIL import Image
from mineru_vl_utils import MinerUClient


# -----------------------------
# 工具：相对 bbox → 像素 bbox
# -----------------------------
def rel_bbox_to_pixel(bbox, width, height):
    x1, y1, x2, y2 = bbox
    return [
        int(x1 * width),
        int(y1 * height),
        int(x2 * width),
        int(y2 * height),
    ]


# -----------------------------
# 工具：截图保存图片
# -----------------------------
def crop_and_save(img: Image.Image, bbox_pixel, save_path: Path, margin=3):
    W, H = img.size
    x1, y1, x2, y2 = bbox_pixel

    log.info(f"[crop_and_save] Image size={img.size}, bbox={bbox_pixel}")

    # --- 1. 越界判断 ---
    if x1 < 0 or y1 < 0 or x2 > W or y2 > H or x2 <= x1 or y2 <= y1:
        log.info("[crop_and_save] BBOX out of range → return original path")
        return None  # 或 return str(original_image_path)

    # --- 2. 判断四条边是否都“紧贴边缘” ---
    touch_left   = x1 <= margin
    touch_top    = y1 <= margin
    touch_right  = x2 >= W - margin
    touch_bottom = y2 >= H - margin

    if touch_left and touch_top and touch_right and touch_bottom:
        log.info("[crop_and_save] BBOX covers almost entire image → skip crop, return original")
        return None  # 或 return str(original_image_path)

    # --- 3. 正常裁剪 ---
    try:
        crop_img = img.crop(bbox_pixel)
        crop_img.save(save_path)
        log.info(f"[crop_and_save] Cropped size={crop_img.size}, saved={save_path}")
        return str(save_path)

    except Exception as e:
        log.info(f"[crop_and_save] ERROR during crop: {e}")
        return None

def transform_sub_bbox(sub_bbox, parent_bbox):
    """把子图内 bbox 映射回原图坐标系"""
    px1, py1, px2, py2 = parent_bbox
    sx1, sy1, sx2, sy2 = sub_bbox

    return [
        px1 + sx1,
        py1 + sy1,
        px1 + sx2,
        py1 + sy2
    ]

# -----------------------------
# 主函数：HTTP异步递归 MinerU
# -----------------------------
async def recursive_run_mineru_http(
    image_path: Path,
    out_dir: Path,
    port:int = 8010,
    max_depth: int = 2,
    current_depth: int = 0,
):
    """
    使用 aio_two_step_extract 执行异步 mineru 提取。
    不依赖中间 JSON 文件。
    自动截图 image/table/list 等元素作为下一轮输入。
    """

    # ---- 深度控制 ----
    if current_depth > max_depth:
        return []

    out_dir.mkdir(parents=True, exist_ok=True)

    log.info(f"[recursive_http] ─ Depth={current_depth}, Image={image_path}")

    # ---- 加载图片 ----
    img = Image.open(image_path)
    W, H = img.size
    log.info(f"[recursive_http] Image size: W={W}, H={H}")

    # ---- 调用 MinerU 异步接口 ----
    try:
        blocks = await run_aio_two_step_extract(image_path, port)
        log.info(f"[recursive_http] MinerU returned {len(blocks)} blocks")
    except Exception as e:
        log.info(f"[recursive_http] MinerU error: {e}")
        return []

    results = []
    sub_images_paths = []  # 下一层递归的图片路径

    # -----------------------------
    # 解析 block 并处理不同类型
    # -----------------------------
    for idx, blk in enumerate(blocks):

        btype = blk.get("type")
        bbox_rel = blk.get("bbox", [0, 0, 1, 1])
        content = blk.get("content")

        log.info(f"  Block[{idx}] type={btype}, bbox_rel={bbox_rel}, "
                 f"content={str(content)[:30] if content else None}")

        bbox_pixel = rel_bbox_to_pixel(bbox_rel, W, H)
        log.info(f"  → bbox_pixel={bbox_pixel}")

        # ---- 文本类 block：直接保存 ----
        if btype in ["title", "text", "paragraph", "caption", "image_caption", "table_caption"]:
            results.append({
                "type": "text",
                "text": content or "",
                "bbox": bbox_pixel,
            })
            log.info(f"    Added TEXT block, content preview: {str(content)[:30]}")

        # ---- 图片类 block：截图作为下一轮输入 ----
        elif btype in ["image", "table", "list"]:

            sub_img_name = f"sub_{current_depth}_{uuid.uuid4()}.png"
            sub_img_path = out_dir / sub_img_name
            
            log.info(f"Try to crop img: {image_path}")
            cropped_path = crop_and_save(img, bbox_pixel, sub_img_path) 
            sub_img_path = cropped_path if cropped_path else image_path
            log.info(f"    Cropped IMAGE block → {sub_img_path}")

            results.append({
                "type": "image",
                "img_path": str(sub_img_path),
                "bbox": bbox_pixel,
            })

            sub_images_paths.append(sub_img_path)

    log.info(f"[recursive_http] Depth={current_depth} → Parsed {len(results)} items, {len(sub_images_paths)} sub-images")

    # ----------------------------------------------
    # 递归处理子图，并将对应的 image 元素替换为子结果
    # ----------------------------------------------
    if current_depth < max_depth and len(sub_images_paths) > 0:

        tasks = [
            recursive_run_mineru_http(
                sub_img_path,
                out_dir,
                port=port,
                max_depth=max_depth,
                current_depth=current_depth + 1
            )
            for sub_img_path in sub_images_paths
        ]

        sub_results_list = await asyncio.gather(*tasks)

        new_results = []

        sub_map = {
            str(path): sub_results_list[i]
            for i, path in enumerate(sub_images_paths)
        }

        for item in results:
            img_path = item.get("img_path")

            # 只对当前层生成的子图做展开替换
            if item["type"] in ["image", 'table', 'list'] and img_path in sub_map:
                parent_bbox = item["bbox"]
                sub_items = sub_map[img_path]

                log.info(f"    Replacing sub-image {img_path} with {len(sub_items)} items")

                for si in sub_items:
                    new_item = si.copy()
                    new_item["bbox"] = transform_sub_bbox(si["bbox"], parent_bbox)
                    new_results.append(new_item)
            else:
                new_results.append(item)

        results = new_results

    return results

# -------------------------------------------------------------------------------------

def get_font_size_for_text(bbox, text, max_font_size=48, min_font_size=10):
    """
    根据文本框的 bbox 和文本长度推算字体大小
    bbox: [xmin, ymin, xmax, ymax]
    text: 要插入的文本
    """
    box_height = bbox[3] - bbox[1]  # 计算文本框高度
    max_chars_per_line = 30         # 最大每行字符数
    lines = ceil(len(text) / max_chars_per_line)  # 计算需要的行数

    font_size = min(box_height // lines, max_font_size)
    return max(font_size, min_font_size)

def generate_ppt_filename(output_path):
    """
    生成一个唯一的 PPT 文件名，基于当前时间戳
    """
    timestamp = time.strftime("%Y%m%d_%H%M%S")
    return f"{output_path}/presentation_{timestamp}.pptx"

def pixels_to_inches(pixels: int, dpi: int = 96) -> float:
    """将像素转换为英寸"""
    return pixels / dpi


def calculate_font_size(text: str, bbox: List[int], text_level: int = None) -> int:
    """
    根据文本框大小、文字内容和文本级别计算合适的字体大小
    """
    # 计算文本框的宽度和高度（像素）
    width = bbox[2] - bbox[0]
    height = bbox[3] - bbox[1]
    
    # 根据文本级别设置基础字体大小
    if text_level == 1:  # 主标题
        base_size = min(height * 0.8, 44)
    elif text_level == 2:  # 子标题
        base_size = min(height * 0.7, 32)
    else:  # 正文
        base_size = min(height * 0.6, 24)
    
    # 根据文本长度调整
    char_count = len(text)
    if char_count > 0:
        chars_per_line = max(1, width / (base_size * 0.6))
        lines_needed = math.ceil(char_count / chars_per_line)
        
        max_lines = max(1, height / (base_size * 1.1))
        if lines_needed > max_lines:
            base_size = base_size * (max_lines / lines_needed)
    
    # 限制字体大小范围
    font_size = max(8, min(base_size, 72))
    
    return int(font_size)


def setup_presentation_size(prs, slide_width_px: int = 1024, slide_height_px: int = 1024):
    """设置PPT尺寸"""
    prs.slide_width = Inches(pixels_to_inches(slide_width_px))
    prs.slide_height = Inches(pixels_to_inches(slide_height_px))
    
    return slide_width_px, slide_height_px


def add_text_element(slide, element: Dict):
    """添加文本元素到幻灯片"""
    bbox = element.get('bbox', [0, 0, 100, 50])
    text = element.get('text', '')
    text_level = element.get('text_level')
    
    # 计算位置和大小
    left = pixels_to_inches(bbox[0])
    top = pixels_to_inches(bbox[1])
    width = pixels_to_inches(bbox[2] - bbox[0])
    height = pixels_to_inches(bbox[3] - bbox[1])
    
    # 计算字体大小
    font_size = calculate_font_size(text, bbox, text_level)
    
    log.info(f"添加文本框:")
    log.info(f"  位置: [{bbox[0]}, {bbox[1]}, {bbox[2]}, {bbox[3]}] 像素")
    log.info(f"  英寸坐标: left={left:.2f}, top={top:.2f}, width={width:.2f}, height={height:.2f}")
    log.info(f"  文本内容: {text[:30]}{'...' if len(text) > 30 else ''}")
    log.info(f"  文本级别: {text_level}, 字体大小: {font_size}pt")
    
    # 添加文本框
    textbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width) * 1.2, Inches(height)
    )
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    # 设置文本内容
    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    
    # 设置字体样式
    paragraph.font.size = Pt(font_size)
    paragraph.font.name = "Comic Sans MS"
    
    # 根据文本级别设置样式
    if text_level == 1:
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.CENTER
        log.info("  样式: 标题(加粗、居中)")
    elif text_level == 2:
        paragraph.font.bold = True
        log.info("  样式: 子标题(加粗)")
    else:
        log.info("  样式: 正文")
    
    return textbox


def add_image_element(slide, element: Dict):
    """添加图片元素到幻灯片"""
    bbox = element.get('bbox', [0, 0, 100, 100])
    img_path = element.get('img_path', '')
    
    # 计算位置和大小
    left = pixels_to_inches(bbox[0])
    top = pixels_to_inches(bbox[1])
    width = pixels_to_inches(bbox[2] - bbox[0])
    height = pixels_to_inches(bbox[3] - bbox[1])
    
    log.info(f"添加图片:")
    log.info(f"  位置: [{bbox[0]}, {bbox[1]}, {bbox[2]}, {bbox[3]}] 像素")
    log.info(f"  英寸坐标: left={left:.2f}, top={top:.2f}, width={width:.2f}, height={height:.2f}")
    log.info(f"  图片路径: {img_path}")
    log.info(f"  图片尺寸: {bbox[2]-bbox[0]}x{bbox[3]-bbox[1]} 像素")
    
    # 检查图片文件是否存在
    if os.path.exists(img_path):
        try:
            log.info("  图片文件存在，正在添加...")
            result = slide.shapes.add_picture(
                img_path,
                Inches(left), Inches(top), Inches(width), Inches(height)
            )
            log.info("  图片添加成功")
            return result
        except Exception as e:
            log.error(f"  添加图片时出错: {e}")
            return add_image_placeholder(slide, bbox, f"Error: {str(e)}")
    else:
        log.warning("  图片文件不存在，使用占位符")
        return add_image_placeholder(slide, bbox, "Image not found")


def add_image_placeholder(slide, bbox: List[int], message: str):
    """添加图片占位符"""
    left = pixels_to_inches(bbox[0])
    top = pixels_to_inches(bbox[1])
    width = pixels_to_inches(bbox[2] - bbox[0])
    height = pixels_to_inches(bbox[3] - bbox[1])
    
    # 添加矩形作为占位符
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
    shape.line.color.rgb = RGBColor(200, 200, 200)
    
    # 添加说明文字
    textbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    text_frame = textbox.text_frame
    text_frame.text = message
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.paragraphs[0].font.size = Pt(10)
    text_frame.paragraphs[0].font.name = "Comic Sans MS"
    text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
    
    return shape

# def robust_parse_json(
#         s: str,
#         merge_dicts: bool = False,          # 想合并时显式打开
#         strip_double_braces: bool = False   # 可选：把 {{ }} → { }
# ) -> Union[Dict[str, Any], List[Any]]:
#     """
#     既能解析普通 JSON，也能从混杂文本中提取多个对象。
    
#     - 支持 // 和 /* */ 注释、尾逗号
#     - 自动去除 Markdown 代码块 ```json … ```、三引号 ''' … '''
#     - 默认返回 dict / list 的原始结构
#     - 若提取到多个独立对象，可用 merge_dicts=True 合并
#     """

#     # ---------- 预处理 ----------
#     # 1) 去掉 ```xxx 代码围栏
#     s = re.sub(r'```[\w]*\s*', '', s)      # 开始围栏
#     s = re.sub(r'```', '', s)              # 结束围栏
#     # 2) 去掉成对的三引号 ''' 或 """
#     s = re.sub(r"^'''|'''$|^\"\"\"|\"\"\"$", '', s.strip())
#     # 3) 可选：{{ }} → { }
#     if strip_double_braces:
#         s = s.replace('{{', '{').replace('}}', '}')
#     # 4) 去注释 + 尾逗号
#     s = _strip_json_comments(s)

#     # ---------- 步骤 1：整体解析 ----------
#     try:
#         return json.loads(s)              
#     except JSONDecodeError:
#         pass                              

#     # ---------- 步骤 2：提取多个对象 ----------
#     objs: List[Any] = _extract_json_objects(s)
#     if not objs:
#         raise ValueError("No valid JSON found.")

#     # 单个对象
#     if len(objs) == 1:
#         return objs[0]

#     # 多对象：根据参数决定
#     if merge_dicts and all(isinstance(o, dict) for o in objs):
#         merged: Dict[str, Any] = {}
#         for o in objs:
#             merged.update(o)
#         return merged
#     return objs

# def _strip_json_comments(s: str) -> str:
#     # s = re.sub(r'/\*.*?\*/', '', s, flags=re.S)          # 块注释
#     # s = re.sub(r'//.*?$',    '', s, flags=re.M)          # 行注释
#     s = re.sub(r',\s*([}\]])', r'\1', s)                 # 尾逗号
#     return s

# def _extract_json_objects(s: str) -> List[Any]:
#     dec = JSONDecoder()
#     idx, n = 0, len(s)
#     objs = []
#     while idx < n:
#         # 找下一个 { 或 [
#         m = re.search(r'[{\[]', s[idx:])
#         if not m:
#             break
#         idx += m.start()
#         try:
#             obj, end = dec.raw_decode(s, idx)
#             objs.append(obj)
#             idx = end
#         except JSONDecodeError:
#             idx += 1
#     return objs


# ========================================================================

#                           For Paper2ExpFigure

# ========================================================================

def pdf_to_pil_images(
    pdf_path: Union[str, Path],
    dpi: int = 300
) -> List[Image.Image]:
    """
    将 PDF 文件的每一页转换为 PIL Image 对象。
    修复了 CMYK/灰度/透明背景导致的白图或花屏问题。
    """
    pdf_path = Path(pdf_path)
    if not pdf_path.exists():
        raise FileNotFoundError(f"PDF 文件不存在: {pdf_path}")

    # 计算缩放比例
    zoom = dpi / 72.0
    matrix = fitz.Matrix(zoom, zoom)

    images: List[Image.Image] = []
    
    # 这里的 flags 用于处理一些复杂的渲染情况（可选，但在某些图表 PDF 中很有用）
    # fitz.pdf.Page.get_pixmap 默认会自动处理大多情况
    
    doc = fitz.open(pdf_path)
    try:
        for page_num in range(len(doc)):
            page = doc[page_num]
            
            # 1. 获取初始 Pixmap
            # alpha=False: 强制背景不透明（默认白色），解决透明底变黑或变白的问题
            pix = page.get_pixmap(matrix=matrix, alpha=False)

            # 2. 关键修复：检查色彩空间并转换
            # pix.n 是通道数。如果不是 3 (RGB)，则强制转换为 RGB
            if pix.n != 3:
                # fitz.csRGB 是 PyMuPDF 内置的 RGB 色彩空间定义
                temp_pix = fitz.Pixmap(fitz.csRGB, pix)
                pix = temp_pix  # 替换为转换后的 RGB pixmap
                # 注意：temp_pix 只是引用，赋值给 pix 后后续逻辑一致，且会自动管理内存

            # 3. 安全转换为 PIL
            # 现在我们可以 100% 确定数据是 RGB 格式的
            img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
            images.append(img)

            log.info(f"[pdf_to_pil_images] 已转换第 {page_num + 1} 页 (模式:{pix.n}通道)，尺寸: {pix.width}x{pix.height}")

    except Exception as e:
        log.error(f"[pdf_to_pil_images] 转换过程出错: {e}")
        raise e
    finally:
        doc.close()

    log.info(f"[pdf_to_pil_images] 完成，共生成 {len(images)} 张图片")
    return images


def _parse_html_table(html_content: str) -> tuple[List[str], List[List[str]]]:
    """
    解析 HTML 表格内容，提取表头和数据行。
    
    参数
    ----
    html_content : str
        HTML 表格字符串，如 <table>...</table>
    
    返回
    ----
    tuple[List[str], List[List[str]]]
        (headers, rows) - 表头列名和数据行
    """
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        log.warning("[_parse_html_table] BeautifulSoup 未安装，使用简单解析")
        return _parse_html_table_simple(html_content)
    
    try:
        soup = BeautifulSoup(html_content, 'html.parser')
        table = soup.find('table')
        
        if not table:
            log.warning("[_parse_html_table] 未找到 table 标签")
            return [], []
        
        # 提取所有行
        rows_data = []
        for tr in table.find_all('tr'):
            row = []
            for cell in tr.find_all(['td', 'th']):
                # 处理 colspan
                colspan = int(cell.get('colspan', 1))
                cell_text = cell.get_text(strip=True)
                row.append(cell_text)
                # 如果有 colspan，添加空单元格
                for _ in range(colspan - 1):
                    row.append('')
            if row:  # 只添加非空行
                rows_data.append(row)
        
        if not rows_data:
            return [], []
        
        # 第一行作为表头
        headers = rows_data[0]
        data_rows = rows_data[1:] if len(rows_data) > 1 else []
        
        return headers, data_rows
        
    except Exception as e:
        log.error(f"[_parse_html_table] 解析失败: {e}")
        return _parse_html_table_simple(html_content)


def _parse_html_table_simple(html_content: str) -> tuple[List[str], List[List[str]]]:
    """
    简单的 HTML 表格解析（不依赖 BeautifulSoup）。
    仅用于备用，可能不够健壮。
    """
    try:
        # 简单正则提取所有 <tr>...</tr> 内容
        import re
        tr_pattern = re.compile(r'<tr>(.*?)</tr>', re.DOTALL | re.IGNORECASE)
        td_pattern = re.compile(r'<t[dh][^>]*>(.*?)</t[dh]>', re.DOTALL | re.IGNORECASE)
        
        rows_data = []
        for tr_match in tr_pattern.finditer(html_content):
            tr_content = tr_match.group(1)
            row = []
            for td_match in td_pattern.finditer(tr_content):
                cell_text = td_match.group(1).strip()
                # 移除 HTML 标签
                cell_text = re.sub(r'<[^>]+>', '', cell_text)
                row.append(cell_text)
            if row:
                rows_data.append(row)
        
        if not rows_data:
            return [], []
        
        headers = rows_data[0]
        data_rows = rows_data[1:] if len(rows_data) > 1 else []
        
        return headers, data_rows
        
    except Exception as e:
        log.error(f"[_parse_html_table_simple] 简单解析失败: {e}")
        return [], []


def extract_tables_from_mineru_results(
    mineru_items: List[Dict[str, Any]],
    min_rows: int = 2,
    min_cols: int = 2,
) -> List[Dict[str, Any]]:
    """
    从 MinerU 识别结果中提取表格数据。

    参数
    ----
    mineru_items : List[Dict[str, Any]]
        MinerU 返回的 items 列表，每个 item 包含 type, bbox, content 等字段
    min_rows : int, default 2
        最小行数，少于此值的表格会被过滤
    min_cols : int, default 2
        最小列数，少于此值的表格会被过滤

    返回
    ----
    List[Dict[str, Any]]
        提取的表格列表，每个表格格式:
        {
            "table_id": str,           # 表格唯一标识
            "headers": List[str],       # 表头列名
            "rows": List[List[str]],    # 数据行
            "caption": str,             # 表格标题/说明
            "bbox": List[int],          # 原始坐标
            "content": str,             # 原始 HTML 内容（如果是 HTML 表格）
        }
    """
    tables = []
    table_idx = 0

    # 用于关联 caption 和 table 的临时存储
    pending_caption = ""

    for item in mineru_items:
        item_type = item.get("type", "")
        content = item.get("content", "")
        bbox = item.get("bbox", [])

        # 处理 table_caption
        if item_type == "table_caption" and content:
            pending_caption = content
            continue

        # 处理表格
        if item_type == "table" and content:
            # MinerU 返回的 content 是 HTML 字符串（如 <table>...</table>）
            headers, rows = _parse_html_table(content)
            
            # 过滤小表格
            if len(headers) < min_cols or len(rows) < min_rows:
                log.debug(f"[extract_tables] 跳过小表格: {len(headers)} 列, {len(rows)} 行")
                continue

            table_id = f"table_{table_idx}"
            table_idx += 1

            tables.append({
                "table_id": table_id,
                "headers": headers,
                "rows": rows,
                "caption": pending_caption,
                "bbox": bbox,
                "content": content,  # 保留原始 HTML
            })

            pending_caption = ""  # 重置
            log.info(f"[extract_tables] 提取表格 {table_id}: {len(headers)} 列, {len(rows)} 行, caption: {pending_caption[:50] if pending_caption else 'N/A'}")

    log.info(f"[extract_tables] 共提取 {len(tables)} 个表格")
    return tables


def extract_text_from_mineru_results(
    mineru_items: List[Dict[str, Any]],
    max_chars: int = 10000,
) -> str:
    """
    从 MinerU 识别结果中提取纯文本内容（用于 paper_idea_extractor）。

    参数
    ----
    mineru_items : List[Dict[str, Any]]
        MinerU 返回的 items 列表
    max_chars : int, default 10000
        最大提取字符数，避免内容过长

    返回
    ----
    str
        提取的文本内容
    """
    text_parts = []
    total_chars = 0

    for item in mineru_items:
        if total_chars >= max_chars:
            break

        item_type = item.get("type", "")
        content = item.get("content", "")

        # 提取文本类型的内容（text, title, table_caption 等）
        if item_type in ["text", "title", "table_caption"] and content:
            text_parts.append(content)
            total_chars += len(content)

    result = "\n\n".join(text_parts)
    
    if total_chars > max_chars:
        result = result[:max_chars] + "..."

    log.info(f"[extract_text] 提取了 {len(text_parts)} 段文本，共 {len(result)} 字符")
    return result


def execute_matplotlib_code(
    code: str,
    output_path: Union[str, Path],
    timeout: int = 30,
    allowed_modules: Optional[List[str]] = None,
) -> Dict[str, Any]:
    """
    安全执行 matplotlib 代码并保存图表。

    参数
    ----
    code : str
        要执行的 matplotlib Python 代码
    output_path : str | Path
        图表输出路径（包含文件名，如 /tmp/chart.png）
    timeout : int, default 30
        执行超时时间（秒）
    allowed_modules : List[str], optional
        允许导入的模块列表，默认为 ["matplotlib", "numpy", "pandas"]

    返回
    ----
    Dict[str, Any]
        {
            "success": bool,
            "output_path": str,      # 成功时返回图片路径
            "error": str,            # 失败时返回错误信息
        }
    """
    import subprocess
    import tempfile

    if allowed_modules is None:
        allowed_modules = ["matplotlib", "numpy", "pandas", "math"]

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    # 安全性检查：禁止危险操作
    dangerous_patterns = [
        r'\bos\.system\b',
        r'\bsubprocess\b',
        r'\beval\b',
        r'\bexec\b',
        r'\bopen\s*\(',
        r'\b__import__\b',
        r'\bimport\s+os\b',
        r'\bimport\s+sys\b',
        r'\bimport\s+subprocess\b',
        r'\bfrom\s+os\b',
        r'\bfrom\s+sys\b',
    ]

    for pattern in dangerous_patterns:
        if re.search(pattern, code):
            return {
                "success": False,
                "output_path": "",
                "error": f"检测到危险操作: {pattern}",
            }

    # 构建完整的执行代码
#     full_code = f'''
# import matplotlib
# matplotlib.use('Agg')  # 非交互式后端
# import matplotlib.pyplot as plt
# import numpy as np

# # 用户代码
# {code}

# # 保存图表
# plt.tight_layout()
# plt.savefig(r"{str(output_path)}", dpi=150, bbox_inches='tight')
# plt.close('all')
# print("SUCCESS")
# '''
    full_code = f'''
{code}
'''

    # 写入临时文件并执行
    try:
        with tempfile.NamedTemporaryFile(mode='w', suffix='.py', delete=False, encoding='utf-8') as f:
            f.write(full_code)
            temp_script = f.name

        result = subprocess.run(
            ['python', temp_script],
            capture_output=True,
            text=True,
            timeout=timeout,
            cwd=str(output_path.parent),
        )

        # 清理临时文件
        try:
            os.unlink(temp_script)
        except Exception:
            pass

        # 判断执行是否成功：返回码为 0 且图片文件存在
        if result.returncode == 0:
            if output_path.exists():
                log.info(f"[execute_matplotlib] 图表生成成功: {output_path}")
                return {
                    "success": True,
                    "output_path": str(output_path),
                    "error": "",
                }
            else:
                log.warning(f"[execute_matplotlib] 代码执行成功但图片未生成")
                return {
                    "success": False,
                    "output_path": "",
                    "error": "代码执行成功但图片未生成",
                }
        else:
            error_msg = result.stderr or result.stdout or "未知错误"
            log.warning(f"[execute_matplotlib] 执行失败 (返回码={result.returncode}): {error_msg}")
            return {
                "success": False,
                "output_path": "",
                "error": error_msg[:500],  # 截断过长的错误信息
            }

    except subprocess.TimeoutExpired:
        log.warning(f"[execute_matplotlib] 执行超时 ({timeout}s)")
        return {
            "success": False,
            "output_path": "",
            "error": f"代码执行超时 ({timeout} 秒)",
        }
    except Exception as e:
        log.error(f"[execute_matplotlib] 执行异常: {e}")
        return {
            "success": False,
            "output_path": "",
            "final_code": full_code,
            "error": str(e),
        }
